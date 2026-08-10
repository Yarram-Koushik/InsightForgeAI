"""
InsightForgeAI – Evidence packs & export helpers

Phase 2.6: evidence packs, CSV, chart HTML/PNG
Phase 4.4: PDF report + PPTX dashboard export

Provides:
  - build_evidence_pack(...)
  - evidence_to_markdown / evidence_to_json
  - dataframe_to_csv_bytes
  - chart_to_html_bytes / chart_to_png_bytes (PNG needs kaleido)
  - build_dashboard_pdf(widgets, title=...) -> ExportPayload
  - build_dashboard_pptx(widgets, title=...) -> ExportPayload
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Optional, Dict, List
from datetime import datetime, timezone
import json
import io

import pandas as pd


MAX_EXPORT_ROWS = 50_000


@dataclass
class ExportPayload:
    filename: str
    mime: str
    data: bytes
    note: Optional[str] = None


def _utc_now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def build_evidence_pack(
    *,
    question: str,
    table_name: str,
    agent_result: Any,
    source_filename: Optional[str] = None,
) -> Dict[str, Any]:
    """Serializable evidence pack for auditability / compliance."""
    result_df = getattr(agent_result, "result_df", None)
    row_count = int(len(result_df)) if result_df is not None else 0
    col_count = int(len(result_df.columns)) if result_df is not None else 0

    pack = {
        "product": "InsightForgeAI",
        "version": "0.4.4",
        "generated_at_utc": _utc_now(),
        "question": question,
        "dataset": {
            "table_name": table_name,
            "source_filename": source_filename,
        },
        "route": {
            "intent": getattr(agent_result, "intent", None),
            "intent_reason": getattr(agent_result, "intent_reason", None),
            "success": bool(getattr(agent_result, "success", False)),
        },
        "sql": getattr(agent_result, "sql", None),
        "insight": getattr(agent_result, "insight", None),
        "message": getattr(agent_result, "message", None),
        "pipeline_steps": list(getattr(agent_result, "steps", []) or []),
        "warnings": list(getattr(agent_result, "warnings", []) or []),
        "error": getattr(agent_result, "error", None),
        "model": {
            "provider": getattr(agent_result, "provider", None),
            "model": getattr(agent_result, "model", None),
        },
        "chart": {
            "type": getattr(agent_result, "chart_type", None),
            "reason": getattr(agent_result, "chart_reason", None),
        },
        "forecast": {
            "method": getattr(agent_result, "forecast_method", None),
            "horizon": getattr(agent_result, "forecast_horizon", None),
            "trend_summary": getattr(agent_result, "trend_summary", None),
            "anomaly_count": len(getattr(agent_result, "anomalies", []) or []),
        },
        "result_shape": {
            "rows": row_count,
            "columns": col_count,
        },
    }
    return pack


def evidence_to_json(pack: Dict[str, Any]) -> bytes:
    return json.dumps(pack, indent=2, default=str).encode("utf-8")


def evidence_to_markdown(pack: Dict[str, Any]) -> bytes:
    lines = [
        f"# InsightForgeAI Evidence Pack",
        f"",
        f"**Generated (UTC):** {pack.get('generated_at_utc')}",
        f"**Question:** {pack.get('question')}",
        f"**Dataset:** `{pack.get('dataset', {}).get('table_name')}`",
        f"**Source file:** {pack.get('dataset', {}).get('source_filename') or '—'}",
        f"",
        f"## Route",
        f"- Intent: **{pack.get('route', {}).get('intent')}**",
        f"- Reason: {pack.get('route', {}).get('intent_reason') or '—'}",
        f"- Success: {pack.get('route', {}).get('success')}",
        f"",
        f"## SQL (evidence)",
        f"```sql",
        f"{pack.get('sql') or '-- none --'}",
        f"```",
        f"",
        f"## Insight",
        f"{pack.get('insight') or pack.get('message') or '—'}",
        f"",
        f"## Pipeline steps",
        f"`{' → '.join(pack.get('pipeline_steps') or [])}`",
        f"",
        f"## Model",
        f"- Provider: {pack.get('model', {}).get('provider') or '—'}",
        f"- Model: {pack.get('model', {}).get('model') or '—'}",
        f"",
        f"## Result shape",
        f"- Rows: {pack.get('result_shape', {}).get('rows')}",
        f"- Columns: {pack.get('result_shape', {}).get('columns')}",
    ]
    warns = pack.get("warnings") or []
    if warns:
        lines += ["", "## Warnings"] + [f"- {w}" for w in warns]
    if pack.get("error"):
        lines += ["", f"**Error:** {pack.get('error')}"]
    return "\n".join(lines).encode("utf-8")


def dataframe_to_csv_bytes(df: Optional[pd.DataFrame], max_rows: int = MAX_EXPORT_ROWS) -> ExportPayload:
    if df is None or not isinstance(df, pd.DataFrame) or df.empty:
        return ExportPayload(filename="empty.csv", mime="text/csv", data=b"", note="No rows to export.")
    note = None
    out = df
    if len(df) > max_rows:
        out = df.head(max_rows)
        note = f"Truncated to {max_rows:,} rows (original {len(df):,})."
    buf = io.StringIO()
    out.to_csv(buf, index=False)
    return ExportPayload(filename="insightforge_result.csv", mime="text/csv", data=buf.getvalue().encode("utf-8"), note=note)


def chart_to_html_bytes(fig: Any, title: str = "chart") -> Optional[ExportPayload]:
    if fig is None:
        return None
    try:
        html = fig.to_html(include_plotlyjs="cdn", full_html=True)
        safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in title)[:40]
        return ExportPayload(filename=f"insightforge_{safe or 'chart'}.html", mime="text/html", data=html.encode("utf-8"))
    except Exception as e:
        return ExportPayload(filename="chart_error.txt", mime="text/plain", data=f"Chart HTML export failed: {e}".encode("utf-8"), note=str(e))


def chart_to_png_bytes(fig: Any, title: str = "chart") -> Optional[ExportPayload]:
    if fig is None:
        return None
    try:
        png = fig.to_image(format="png", scale=2)
        safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in title)[:40]
        return ExportPayload(filename=f"insightforge_{safe or 'chart'}.png", mime="image/png", data=png)
    except Exception:
        return None


def safe_filename_part(text: str, max_len: int = 32) -> str:
    raw = (text or "export").strip().lower().replace(" ", "_")
    cleaned = "".join(c if c.isalnum() or c in "-_" else "_" for c in raw)
    return (cleaned[:max_len] or "export")


# ---------------------------------------------------------------------------
# Phase 4.4 – Dashboard PDF / PPTX
# ---------------------------------------------------------------------------

def _widget_as_dict(w: Any) -> Dict[str, Any]:
    if hasattr(w, "to_dict"):
        return w.to_dict()
    if isinstance(w, dict):
        return w
    return {
        "title": getattr(w, "title", None) or getattr(w, "question", "Widget"),
        "question": getattr(w, "question", ""),
        "table_name": getattr(w, "table_name", ""),
        "sql": getattr(w, "sql", None),
        "insight": getattr(w, "insight", None),
        "chart_type": getattr(w, "chart_type", None),
        "grounding_line": getattr(w, "grounding_line", None),
        "status": getattr(w, "status", "ok"),
        "error": getattr(w, "error", None),
        "last_row_count": getattr(w, "last_row_count", None),
        "created_at": getattr(w, "created_at", None),
    }


def build_dashboard_pdf(
    widgets: List[Any],
    *,
    title: str = "InsightForgeAI Dashboard",
    subtitle: Optional[str] = None,
) -> ExportPayload:
    """
    Build a multi-page PDF from pinned dashboard widgets.
    Uses reportlab. Charts are described textually (PNG embedding optional via kaleido).
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted, HRFlowable
        from reportlab.lib.colors import HexColor
    except ImportError as e:
        return ExportPayload(
            filename="dashboard_error.txt",
            mime="text/plain",
            data=f"PDF export requires reportlab: pip install reportlab\n{e}".encode("utf-8"),
            note="Missing reportlab",
        )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=16 * mm,
        bottomMargin=16 * mm,
        title=title,
        author="InsightForgeAI",
    )
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="IFTitle", parent=styles["Heading1"], fontSize=18, spaceAfter=6))
    styles.add(ParagraphStyle(name="IFSub", parent=styles["Normal"], fontSize=9, textColor=HexColor("#555555"), spaceAfter=12))
    styles.add(ParagraphStyle(name="IFH2", parent=styles["Heading2"], fontSize=13, spaceBefore=8, spaceAfter=4))
    styles.add(ParagraphStyle(name="IFBody", parent=styles["Normal"], fontSize=10, leading=13, spaceAfter=6))
    styles.add(ParagraphStyle(name="IFMeta", parent=styles["Normal"], fontSize=8, textColor=HexColor("#666666"), spaceAfter=4))
    styles.add(ParagraphStyle(name="IFSQL", parent=styles["Code"], fontSize=7, leading=9, leftIndent=4, spaceAfter=8))

    story = []
    story.append(Paragraph(title.replace("&", "&"), styles["IFTitle"]))
    sub = subtitle or f"Generated {_utc_now()} · {len(widgets)} widget(s)"
    story.append(Paragraph(sub.replace("&", "&"), styles["IFSub"]))
    story.append(HRFlowable(width="100%", thickness=0.5, color=HexColor("#cccccc"), spaceAfter=10))

    if not widgets:
        story.append(Paragraph("No widgets pinned yet.", styles["IFBody"]))
    else:
        for i, raw in enumerate(widgets):
            w = _widget_as_dict(raw)
            wtitle = (w.get("title") or w.get("question") or f"Widget {i+1}").replace("&", "&")
            story.append(Paragraph(f"{i+1}. {wtitle}", styles["IFH2"]))

            meta_bits = []
            if w.get("table_name"):
                meta_bits.append(f"Dataset: <b>{w['table_name']}</b>")
            if w.get("status"):
                meta_bits.append(f"Status: {w['status']}")
            if w.get("chart_type"):
                meta_bits.append(f"Chart: {w['chart_type']}")
            if w.get("last_row_count") is not None:
                meta_bits.append(f"Rows: {w['last_row_count']}")
            if meta_bits:
                story.append(Paragraph(" · ".join(meta_bits), styles["IFMeta"]))

            if w.get("question"):
                story.append(Paragraph(f"<b>Question:</b> {str(w['question']).replace('&', '&')}", styles["IFBody"]))

            if w.get("grounding_line"):
                story.append(Paragraph(f"<i>Grounding: {str(w['grounding_line']).replace('&', '&')}</i>", styles["IFMeta"]))

            if w.get("insight"):
                insight_txt = str(w["insight"]).replace("&", "&").replace("\n", "<br/>")
                story.append(Paragraph(f"<b>Insight</b><br/>{insight_txt}", styles["IFBody"]))

            if w.get("sql"):
                sql_safe = str(w["sql"])[:4000]
                story.append(Paragraph("<b>SQL (evidence)</b>", styles["IFMeta"]))
                story.append(Preformatted(sql_safe, styles["IFSQL"]))

            if w.get("error"):
                story.append(Paragraph(f"<font color='red'><b>Error:</b> {str(w['error']).replace('&', '&')}</font>", styles["IFBody"]))

            if i < len(widgets) - 1:
                story.append(Spacer(1, 6))
                story.append(HRFlowable(width="100%", thickness=0.3, color=HexColor("#eeeeee"), spaceAfter=6))

    try:
        doc.build(story)
        data = buf.getvalue()
        fname = f"insightforge_dashboard_{safe_filename_part(title)}.pdf"
        return ExportPayload(filename=fname, mime="application/pdf", data=data)
    except Exception as e:
        return ExportPayload(
            filename="dashboard_pdf_error.txt",
            mime="text/plain",
            data=f"PDF build failed: {e}".encode("utf-8"),
            note=str(e),
        )


def build_dashboard_pptx(
    widgets: List[Any],
    *,
    title: str = "InsightForgeAI Dashboard",
    subtitle: Optional[str] = None,
) -> ExportPayload:
    """
    Build a PowerPoint deck: title slide + one slide per widget.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
    except ImportError as e:
        return ExportPayload(
            filename="dashboard_error.txt",
            mime="text/plain",
            data=f"PPTX export requires python-pptx: pip install python-pptx\n{e}".encode("utf-8"),
            note="Missing python-pptx",
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x1A, 0x1A, 0x2E)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.6))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle or f"Generated {_utc_now()} · {len(widgets)} widget(s)"
    sp.font.size = Pt(14)
    sp.font.color.rgb = RgbColor(0x55, 0x55, 0x55)

    if not widgets:
        slide2 = prs.slides.add_slide(blank)
        t = slide2.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11), Inches(1))
        t.text_frame.paragraphs[0].text = "No widgets pinned yet."
        t.text_frame.paragraphs[0].font.size = Pt(18)

    for i, raw in enumerate(widgets):
        w = _widget_as_dict(raw)
        slide = prs.slides.add_slide(blank)

        h = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        hp = h.text_frame.paragraphs[0]
        hp.text = f"{i+1}. {w.get('title') or w.get('question') or 'Widget'}"
        hp.font.size = Pt(20)
        hp.font.bold = True
        hp.font.color.rgb = RgbColor(0x1A, 0x1A, 0x2E)

        y = 1.0
        meta = []
        if w.get("table_name"):
            meta.append(f"Dataset: {w['table_name']}")
        if w.get("status"):
            meta.append(f"Status: {w['status']}")
        if w.get("chart_type"):
            meta.append(f"Chart: {w['chart_type']}")
        if meta:
            mbox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.35))
            mbox.text_frame.paragraphs[0].text = " · ".join(meta)
            mbox.text_frame.paragraphs[0].font.size = Pt(11)
            mbox.text_frame.paragraphs[0].font.color.rgb = RgbColor(0x66, 0x66, 0x66)
            y += 0.4

        if w.get("question"):
            qbox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.5))
            qbox.text_frame.paragraphs[0].text = f"Q: {w['question']}"
            qbox.text_frame.paragraphs[0].font.size = Pt(12)
            y += 0.55

        if w.get("insight"):
            ibox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(2.2))
            tf = ibox.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = "Insight"
            p0.font.bold = True
            p0.font.size = Pt(12)
            p1 = tf.add_paragraph()
            p1.text = str(w["insight"])[:1500]
            p1.font.size = Pt(11)
            y += 2.4

        if w.get("sql"):
            sbox = slide.shapes.add_textbox(Inches(0.5), Inches(min(y, 5.5)), Inches(12.3), Inches(1.5))
            tf = sbox.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = "SQL"
            p0.font.bold = True
            p0.font.size = Pt(10)
            p1 = tf.add_paragraph()
            p1.text = str(w["sql"])[:1200]
            p1.font.size = Pt(8)
            p1.font.name = "Courier New"

        if w.get("error"):
            ebox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4))
            ebox.text_frame.paragraphs[0].text = f"Error: {w['error']}"
            ebox.text_frame.paragraphs[0].font.size = Pt(10)
            ebox.text_frame.paragraphs[0].font.color.rgb = RgbColor(0xC0, 0x00, 0x00)

    buf = io.BytesIO()
    try:
        prs.save(buf)
        data = buf.getvalue()
        fname = f"insightforge_dashboard_{safe_filename_part(title)}.pptx"
        return ExportPayload(
            filename=fname,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            data=data,
        )
    except Exception as e:
        return ExportPayload(
            filename="dashboard_pptx_error.txt",
            mime="text/plain",
            data=f"PPTX build failed: {e}".encode("utf-8"),
            note=str(e),
        )
